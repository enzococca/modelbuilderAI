"""File search tool — search files across local PC, Dropbox, Google Drive, OneDrive."""

from __future__ import annotations

import asyncio
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from tools import BaseTool


class FileSearchTool(BaseTool):
    """Search for files across multiple backends."""

    name = "file_search"
    description = (
        "Search files across local filesystem, Dropbox, Google Drive, "
        "and Microsoft OneDrive. Returns file names, paths, sizes, dates."
    )

    async def execute(self, input_text: str, **kwargs: Any) -> str:
        source = kwargs.get("source", "local")
        mode = kwargs.get("mode", "filename")  # "filename" or "content"
        max_results = int(kwargs.get("max_results", 20))

        # Remove keys already extracted to avoid duplicate keyword args
        fwd = {k: v for k, v in kwargs.items() if k not in ("source", "mode", "max_results")}

        # Content search mode — search INSIDE documents
        if mode == "content" and source == "local":
            return await self._search_local_content(input_text, max_results=max_results, **fwd)

        dispatch = {
            "local": self._search_local,
            "dropbox": self._search_dropbox,
            "gdrive": self._search_gdrive,
            "onedrive": self._search_onedrive,
        }

        handler = dispatch.get(source)
        if handler is None:
            return f"Unknown source: {source}. Available: {', '.join(dispatch.keys())}"

        try:
            return await handler(input_text, max_results=max_results, **fwd)
        except ImportError as e:
            return f"Missing dependency for {source}: {e}. Install it with pip."
        except Exception as e:
            return f"File search error ({source}): {e}"

    # ── local ──────────────────────────────────────────────────

    async def _search_local(self, query: str, **kwargs: Any) -> str:
        """Search local filesystem using filename matching."""
        from config import settings

        max_results = int(kwargs.get("max_results", 20))
        roots_str = kwargs.get("roots", "") or settings.file_search_local_roots

        if not roots_str:
            roots = [str(Path.home())]
        else:
            roots = [r.strip() for r in roots_str.split(",") if r.strip()]

        # Run blocking I/O in thread
        results = await asyncio.to_thread(
            self._local_search_sync, query, roots, max_results
        )
        return self._format_results("Local", query, results)

    @staticmethod
    def _local_search_sync(
        query: str, roots: list[str], max_results: int
    ) -> list[dict[str, Any]]:
        """Synchronous local file search (runs in thread)."""
        results: list[dict[str, Any]] = []
        q = query.lower()

        for root in roots:
            root_path = Path(root)
            if not root_path.exists():
                continue
            try:
                for p in root_path.rglob("*"):
                    if len(results) >= max_results:
                        break
                    if not p.is_file():
                        continue
                    if q in p.name.lower():
                        try:
                            stat = p.stat()
                            results.append({
                                "name": p.name,
                                "path": str(p),
                                "size": stat.st_size,
                                "modified": datetime.fromtimestamp(
                                    stat.st_mtime, tz=timezone.utc
                                ).isoformat(),
                            })
                        except OSError:
                            continue
            except PermissionError:
                continue
            if len(results) >= max_results:
                break

        return results

    # ── local content search ─────────────────────────────────

    async def _search_local_content(self, query: str, **kwargs: Any) -> str:
        """Search INSIDE local files (PDF, DOCX, PPTX, TXT, etc.) for a query string.

        Recursively walks all subdirectories, extracts text from supported
        file types, and returns files + relevant excerpts where the query appears.
        """
        from config import settings

        max_results = int(kwargs.get("max_results", 20))
        roots_str = kwargs.get("roots", "") or settings.file_search_local_roots
        extensions = kwargs.get("extensions", "pdf,docx,pptx,txt,md,csv")
        ext_set = {
            f".{e.strip().lower().lstrip('.')}"
            for e in extensions.split(",")
            if e.strip()
        }

        if not roots_str:
            roots = [str(Path.home())]
        else:
            roots = [r.strip() for r in roots_str.split(",") if r.strip()]

        results = await asyncio.to_thread(
            self._content_search_sync, query, roots, ext_set, max_results
        )

        if not results:
            return f"No files containing '{query}' found in: {', '.join(roots)}"

        lines = [
            f"**Content Search Results** — query: *{query}*\n",
            f"Searched extensions: {extensions}\n",
        ]
        for i, r in enumerate(results, 1):
            lines.append(f"### {i}. {r['name']}")
            lines.append(f"Path: `{r['path']}`\n")
            for excerpt in r["excerpts"]:
                lines.append(f"> ...{excerpt}...\n")
            lines.append("")

        lines.append(f"**Total: {len(results)} file(s) containing '{query}'**")

        # Also append full extracted text for downstream agents
        lines.append("\n\n---\n\n## Full Extracted Text\n")
        for r in results:
            lines.append(f"### {r['name']}\n")
            lines.append(r.get("full_text", "")[:8000])
            lines.append("\n---\n")

        return "\n".join(lines)

    @staticmethod
    def _content_search_sync(
        query: str,
        roots: list[str],
        ext_set: set[str],
        max_results: int,
    ) -> list[dict[str, Any]]:
        """Walk directories, extract text from documents, search for query."""
        q_lower = query.lower()
        results: list[dict[str, Any]] = []

        for root in roots:
            root_path = Path(root)
            if not root_path.exists():
                continue
            try:
                for p in root_path.rglob("*"):
                    if len(results) >= max_results:
                        break
                    if not p.is_file():
                        continue
                    if p.suffix.lower() not in ext_set:
                        continue
                    try:
                        text = _extract_text(p)
                        if not text:
                            continue
                        if q_lower not in text.lower():
                            continue
                        # Found a match — extract relevant excerpts
                        excerpts = _extract_excerpts(text, query, context_chars=200)
                        results.append({
                            "name": p.name,
                            "path": str(p),
                            "excerpts": excerpts[:5],  # max 5 excerpts per file
                            "full_text": text,
                        })
                    except Exception:
                        continue
            except PermissionError:
                continue
            if len(results) >= max_results:
                break

        return results

    # ── dropbox ────────────────────────────────────────────────

    async def _search_dropbox(self, query: str, **kwargs: Any) -> str:
        """Search Dropbox using files_search_v2."""
        import dropbox
        from config import settings

        if not settings.dropbox_refresh_token:
            return (
                "Dropbox not configured. Set dropbox_app_key, "
                "dropbox_app_secret, and dropbox_refresh_token in Settings."
            )

        max_results = int(kwargs.get("max_results", 20))

        # Run blocking SDK call in thread
        def _do_search() -> list[dict[str, Any]]:
            dbx = dropbox.Dropbox(
                app_key=settings.dropbox_app_key,
                app_secret=settings.dropbox_app_secret,
                oauth2_refresh_token=settings.dropbox_refresh_token,
            )
            result = dbx.files_search_v2(query)
            items: list[dict[str, Any]] = []
            for match in result.matches[:max_results]:
                metadata = match.metadata.get_metadata()
                items.append({
                    "name": metadata.name,
                    "path": getattr(metadata, "path_display", ""),
                    "size": getattr(metadata, "size", 0),
                    "modified": str(getattr(metadata, "server_modified", "")),
                })
            return items

        results = await asyncio.to_thread(_do_search)
        return self._format_results("Dropbox", query, results)

    # ── gdrive ─────────────────────────────────────────────────

    async def _search_gdrive(self, query: str, **kwargs: Any) -> str:
        """Search Google Drive via API."""
        from config import settings

        if not settings.google_drive_credentials_json:
            return (
                "Google Drive not configured. Set google_drive_credentials_json "
                "(path to service account JSON) in Settings."
            )

        max_results = int(kwargs.get("max_results", 20))

        def _do_search() -> list[dict[str, Any]]:
            from google.oauth2 import service_account
            from googleapiclient.discovery import build

            creds = service_account.Credentials.from_service_account_file(
                settings.google_drive_credentials_json,
                scopes=["https://www.googleapis.com/auth/drive.readonly"],
            )
            if settings.google_drive_delegated_user:
                creds = creds.with_subject(settings.google_drive_delegated_user)

            service = build("drive", "v3", credentials=creds)

            # Escape single quotes in query
            safe_q = query.replace("'", "\\'")
            response = (
                service.files()
                .list(
                    q=f"name contains '{safe_q}' and trashed=false",
                    pageSize=max_results,
                    fields="files(id, name, mimeType, size, modifiedTime, webViewLink)",
                )
                .execute()
            )

            items: list[dict[str, Any]] = []
            for f in response.get("files", []):
                items.append({
                    "name": f["name"],
                    "path": f.get("webViewLink", f["id"]),
                    "size": int(f.get("size", 0)),
                    "modified": f.get("modifiedTime", ""),
                    "mime_type": f.get("mimeType", ""),
                })
            return items

        results = await asyncio.to_thread(_do_search)
        return self._format_results("Google Drive", query, results)

    # ── onedrive ───────────────────────────────────────────────

    async def _search_onedrive(self, query: str, **kwargs: Any) -> str:
        """Search OneDrive via Microsoft Graph REST API."""
        from config import settings

        if not settings.microsoft_client_id or not settings.microsoft_tenant_id:
            return (
                "Microsoft not configured. Set microsoft_tenant_id, "
                "microsoft_client_id, and microsoft_client_secret in Settings."
            )

        max_results = int(kwargs.get("max_results", 20))

        token = await _get_ms_token(settings)
        headers = {"Authorization": f"Bearer {token}"}
        user_id = settings.microsoft_user_id or "me"

        import httpx

        async with httpx.AsyncClient(timeout=15) as client:
            url = (
                f"https://graph.microsoft.com/v1.0/users/{user_id}"
                f"/drive/root/search(q='{query}')"
            )
            resp = await client.get(url, headers=headers, params={"$top": max_results})
            resp.raise_for_status()
            data = resp.json()

        results: list[dict[str, Any]] = []
        for item in data.get("value", []):
            results.append({
                "name": item.get("name", ""),
                "path": item.get("webUrl", ""),
                "size": item.get("size", 0),
                "modified": item.get("lastModifiedDateTime", ""),
            })

        return self._format_results("OneDrive", query, results)

    # ── formatting ─────────────────────────────────────────────

    @staticmethod
    def _format_results(
        source: str, query: str, results: list[dict[str, Any]]
    ) -> str:
        if not results:
            return f"No files found on {source} for: {query}"
        lines = [f"**File Search Results** ({source}) — query: *{query}*\n"]
        for i, r in enumerate(results, 1):
            size_kb = r.get("size", 0) / 1024
            lines.append(f"{i}. **{r['name']}**")
            lines.append(f"   Path: {r['path']}")
            lines.append(
                f"   Size: {size_kb:.1f} KB | Modified: {r.get('modified', 'N/A')}"
            )
            if r.get("mime_type"):
                lines.append(f"   Type: {r['mime_type']}")
            lines.append("")
        lines.append(f"Total: {len(results)} file(s)")
        return "\n".join(lines)


# ── shared Microsoft token helper ──────────────────────────────

# ── document text extraction helpers ───────────────────────────


def _extract_text(path: Path) -> str:
    """Extract text content from a file. Supports PDF, DOCX, PPTX, and text files."""
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(path)
    elif suffix == ".docx":
        return _extract_docx(path)
    elif suffix == ".pptx":
        return _extract_pptx(path)
    elif suffix in (".txt", ".md", ".csv", ".py", ".js", ".ts", ".html", ".xml", ".yaml", ".yml"):
        try:
            return path.read_text(errors="replace")
        except Exception:
            return ""
    return ""


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text)
        return "\n\n".join(pages)
    except Exception:
        return ""


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    except Exception:
        return ""


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX (PowerPoint)."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                texts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)
    except Exception:
        return ""


def _extract_excerpts(text: str, query: str, context_chars: int = 200) -> list[str]:
    """Extract text excerpts around each occurrence of query."""
    excerpts = []
    text_lower = text.lower()
    q_lower = query.lower()
    start = 0

    while True:
        idx = text_lower.find(q_lower, start)
        if idx == -1:
            break
        excerpt_start = max(0, idx - context_chars)
        excerpt_end = min(len(text), idx + len(query) + context_chars)
        excerpt = text[excerpt_start:excerpt_end].replace("\n", " ").strip()
        excerpts.append(excerpt)
        start = idx + len(query)
        if len(excerpts) >= 10:
            break

    return excerpts


async def _get_ms_token(settings: Any) -> str:
    """Obtain Microsoft Graph access token via client credentials flow."""
    import httpx

    url = (
        f"https://login.microsoftonline.com/"
        f"{settings.microsoft_tenant_id}/oauth2/v2.0/token"
    )
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            url,
            data={
                "grant_type": "client_credentials",
                "client_id": settings.microsoft_client_id,
                "client_secret": settings.microsoft_client_secret,
                "scope": "https://graph.microsoft.com/.default",
            },
        )
        resp.raise_for_status()
        return resp.json()["access_token"]
